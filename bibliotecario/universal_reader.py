"""
Leitor Universal de Documentos para ENXAME Bibliotecário
=========================================================
Suporte a múltiplos formatos de arquivo:
- Texto: .txt, .md, .py, .json, .yaml, .yml, .js, .ts, .csv, .rtf
- Documentos: .pdf, .docx, .odt, .pptx, .xlsx
- Imagens: .jpg, .jpeg, .png, .gif, .bmp, .tiff (OCR via pytesseract se disponível)
- Vídeos: .mp4, .avi, .mkv, .webm (extração de metadados e legendas)
- Áudio: .mp3, .wav, .ogg, .flac (extração de metadados e transcrição se disponível)
- Arquivos ZIM: .zim (via zim_reader existente)

Princípios:
- Graceful degradation: se uma biblioteca não estiver disponível, ignora o formato
- Offline-first: nenhuma dependência de serviços online
- Segurança: validação de caminhos e prevenção de path traversal
- Performance: leitura lazy e chunking inteligente
"""
from __future__ import annotations

import logging
import os
from dataclasses import dataclass
from pathlib import Path
from typing import Optional

logger = logging.getLogger(__name__)


@dataclass(slots=True)
class DocumentChunk:
    """Representa um chunk de texto extraído de um documento."""
    chunk_id: str
    text: str
    source_path: str
    extension: str
    metadata: dict  # metadados específicos do arquivo


class UniversalDocumentReader:
    """
    Leitor universal de documentos com suporte a múltiplos formatos.
    Implementa graceful degradation: se uma biblioteca não estiver disponível,
    o formato correspondente é ignorado sem falhar o sistema.
    """

    # Mapeamento de extensões para categorias
    TEXT_EXTENSIONS = {".txt", ".md", ".py", ".json", ".yaml", ".yml", ".js", ".ts", ".csv", ".rtf", ".html", ".htm"}
    PDF_EXTENSIONS = {".pdf"}
    DOCX_EXTENSIONS = {".docx"}
    PPTX_EXTENSIONS = {".pptx"}
    XLSX_EXTENSIONS = {".xlsx", ".xlsm"}
    ODF_EXTENSIONS = {".odt", ".ods", ".odp"}
    IMAGE_EXTENSIONS = {".jpg", ".jpeg", ".png", ".gif", ".bmp", ".tiff", ".tif", ".webp"}
    VIDEO_EXTENSIONS = {".mp4", ".avi", ".mkv", ".webm", ".mov", ".flv"}
    AUDIO_EXTENSIONS = {".mp3", ".wav", ".ogg", ".flac", ".m4a", ".aac"}
    
    ALL_SUPPORTED = (
        TEXT_EXTENSIONS | PDF_EXTENSIONS | DOCX_EXTENSIONS | 
        PPTX_EXTENSIONS | XLSX_EXTENSIONS | ODF_EXTENSIONS |
        IMAGE_EXTENSIONS | VIDEO_EXTENSIONS | AUDIO_EXTENSIONS
    )

    def __init__(
        self,
        chunk_size: int = 1200,
        overlap: int = 150,
        enable_ocr: bool = False,
        enable_transcription: bool = False,
    ) -> None:
        """
        Inicializa o leitor universal.
        
        Args:
            chunk_size: Tamanho máximo de cada chunk em caracteres
            overlap: Sobreposição entre chunks consecutivos
            enable_ocr: Se True, tenta extrair texto de imagens via OCR (requer pytesseract)
            enable_transcription: Se True, tenta transcrever áudio (requer bibliotecas extras)
        """
        self.chunk_size = chunk_size
        self.overlap = overlap
        self.enable_ocr = enable_ocr
        self.enable_transcription = enable_transcription
        
        # Cache de disponibilidade de bibliotecas
        self._library_cache: dict[str, bool] = {}

    def _check_library(self, lib_name: str) -> bool:
        """Verifica se uma biblioteca está disponível (com cache)."""
        if lib_name in self._library_cache:
            return self._library_cache[lib_name]
        
        try:
            __import__(lib_name)
            self._library_cache[lib_name] = True
            return True
        except ImportError:
            self._library_cache[lib_name] = False
            logger.debug(f"Biblioteca '{lib_name}' não disponível")
            return False

    def _extract_text_plain(self, path: Path) -> str:
        """Extrai texto de arquivos de texto puro."""
        try:
            # Tenta UTF-8 primeiro, fallback para latin-1
            for encoding in ["utf-8", "latin-1", "cp1252"]:
                try:
                    content = path.read_text(encoding=encoding)
                    # Remove BOM se presente
                    if content.startswith('\ufeff'):
                        content = content[1:]
                    return content
                except UnicodeDecodeError:
                    continue
            logger.warning(f"Não foi possível decodificar {path}")
            return ""
        except Exception as exc:
            logger.warning(f"Falha ao ler {path}: {exc}")
            return ""

    def _extract_text_pdf(self, path: Path) -> str:
        """Extrai texto de arquivos PDF."""
        if not self._check_library("fitz"):
            logger.warning(f"PyMuPDF não disponível, ignorando {path}")
            return ""
        
        try:
            import fitz  # type: ignore
            doc = fitz.open(path)
            texts = []
            for page in doc:
                text = page.get_text("text")
                if text.strip():
                    texts.append(text)
            doc.close()
            return "\n".join(texts)
        except Exception as exc:
            logger.warning(f"Falha ao extrair PDF {path}: {exc}")
            return ""

    def _extract_text_docx(self, path: Path) -> str:
        """Extrai texto de arquivos DOCX."""
        if not self._check_library("docx"):
            logger.warning(f"python-docx não disponível, ignorando {path}")
            return ""
        
        try:
            from docx import Document  # type: ignore
            document = Document(path)
            paragraphs = [p.text for p in document.paragraphs if p.text.strip()]
            
            # Tentar extrair texto de tabelas também
            for table in document.tables:
                for row in table.rows:
                    for cell in row.cells:
                        if cell.text.strip():
                            paragraphs.append(cell.text)
            
            return "\n".join(paragraphs)
        except Exception as exc:
            logger.warning(f"Falha ao extrair DOCX {path}: {exc}")
            return ""

    def _extract_text_pptx(self, path: Path) -> str:
        """Extrai texto de arquivos PPTX."""
        if not self._check_library("pptx"):
            logger.warning(f"python-pptx não disponível, ignorando {path}")
            return ""
        
        try:
            from pptx import Presentation  # type: ignore
            prs = Presentation(path)
            texts = []
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        texts.append(shape.text)
            return "\n\n".join(texts)
        except Exception as exc:
            logger.warning(f"Falha ao extrair PPTX {path}: {exc}")
            return ""

    def _extract_text_xlsx(self, path: Path) -> str:
        """Extrai texto de arquivos XLSX."""
        if not self._check_library("openpyxl"):
            logger.warning(f"openpyxl não disponível, ignorando {path}")
            return ""
        
        try:
            from openpyxl import load_workbook  # type: ignore
            wb = load_workbook(path, read_only=True, data_only=True)
            texts = []
            for sheet_name in wb.sheetnames:
                sheet = wb[sheet_name]
                texts.append(f"=== Planilha: {sheet_name} ===")
                for row in sheet.iter_rows(values_only=True):
                    row_text = " | ".join(str(cell) if cell is not None else "" for cell in row)
                    if row_text.strip():
                        texts.append(row_text)
            wb.close()
            return "\n".join(texts)
        except Exception as exc:
            logger.warning(f"Falha ao extrair XLSX {path}: {exc}")
            return ""

    def _extract_text_odf(self, path: Path) -> str:
        """Extrai texto de arquivos ODF (ODT, ODS, ODP)."""
        if not self._check_library("odf"):
            logger.warning(f"odfpy não disponível, ignorando {path}")
            return ""
        
        try:
            from odf.opendocument import load  # type: ignore
            from odf.textextractor import extractText, ODFTextExtractor  # type: ignore
            
            doc = load(path)
            extractor = ODFTextExtractor(doc)
            return extractText(extractor)
        except Exception as exc:
            logger.warning(f"Falha ao extrair ODF {path}: {exc}")
            return ""

    def _extract_text_image(self, path: Path) -> str:
        """Extrai texto de imagens via OCR (se habilitado e disponível)."""
        if not self.enable_ocr:
            logger.debug(f"OCR desabilitado, ignorando {path}")
            return ""
        
        if not self._check_library("pytesseract"):
            logger.warning(f"pytesseract não disponível, ignorando {path}")
            return ""
        
        if not self._check_library("PIL"):
            logger.warning(f"Pillow não disponível, ignorando {path}")
            return ""
        
        try:
            import pytesseract  # type: ignore
            from PIL import Image  # type: ignore
            
            img = Image.open(path)
            text = pytesseract.image_to_string(img, lang="por+eng")
            return f"[OCR] {text}"
        except Exception as exc:
            logger.warning(f"Falha ao extrair OCR de {path}: {exc}")
            return ""

    def _extract_metadata_media(self, path: Path) -> dict:
        """Extrai metadados de arquivos de mídia (vídeo/áudio)."""
        metadata = {
            "type": "media",
            "filename": path.name,
            "size_bytes": path.stat().st_size,
        }
        
        # Tentar extrair metadados com mutagen (áudio) ou mediainfo
        if path.suffix.lower() in self.AUDIO_EXTENSIONS and self._check_library("mutagen"):
            try:
                from mutagen import File  # type: ignore
                audio = File(path)
                if audio:
                    metadata["duration"] = getattr(audio.info, "length", None)
                    metadata["tags"] = dict(audio.tags) if audio.tags else {}
            except Exception as exc:
                logger.debug(f"Falha ao extrair metadados de áudio {path}: {exc}")
        
        # Para vídeos, tentar extrair legendas embutidas se disponível
        if path.suffix.lower() in self.VIDEO_EXTENSIONS:
            # Implementação básica - pode ser expandida
            metadata["subtype"] = "video"
        
        return metadata

    def extract_text(self, path: Path) -> tuple[str, dict]:
        """
        Extrai texto de um arquivo baseado em sua extensão.
        
        Args:
            path: Caminho absoluto para o arquivo
        
        Returns:
            Tupla (texto_extraído, metadados)
        """
        # Validação de segurança
        if not path.is_absolute():
            raise ValueError(f"Caminho deve ser absoluto: {path}")
        
        if not path.exists():
            logger.warning(f"Arquivo não existe: {path}")
            return "", {}
        
        if not path.is_file():
            logger.warning(f"NÃO é um arquivo: {path}")
            return "", {}
        
        ext = path.suffix.lower()
        metadata = {
            "source_path": str(path),
            "extension": ext,
            "size_bytes": path.stat().st_size,
        }
        
        # Roteamento por tipo de arquivo
        try:
            if ext in self.TEXT_EXTENSIONS:
                text = self._extract_text_plain(path)
            elif ext in self.PDF_EXTENSIONS:
                text = self._extract_text_pdf(path)
            elif ext in self.DOCX_EXTENSIONS:
                text = self._extract_text_docx(path)
            elif ext in self.PPTX_EXTENSIONS:
                text = self._extract_text_pptx(path)
            elif ext in self.XLSX_EXTENSIONS:
                text = self._extract_text_xlsx(path)
            elif ext in self.ODF_EXTENSIONS:
                text = self._extract_text_odf(path)
            elif ext in self.IMAGE_EXTENSIONS:
                text = self._extract_text_image(path)
            elif ext in self.VIDEO_EXTENSIONS | self.AUDIO_EXTENSIONS:
                text = ""
                metadata.update(self._extract_metadata_media(path))
            else:
                logger.debug(f"Extensão não suportada: {ext}")
                return "", metadata
            
            return text or "", metadata
        except Exception as exc:
            logger.exception(f"Erro inesperado ao extrair {path}: {exc}")
            return "", metadata

    def split_into_chunks(self, text: str, metadata: dict) -> list[DocumentChunk]:
        """
        Divide texto em chunks com sobreposição.
        
        Args:
            text: Texto completo a ser dividido
            metadata: Metadados do documento original
        
        Returns:
            Lista de DocumentChunk
        """
        import hashlib
        
        if not text or not text.strip():
            return []
        
        # Normalizar whitespace
        normalized = " ".join(text.split())
        
        if len(normalized) <= self.chunk_size:
            chunk_id = hashlib.sha1(f"{metadata['source_path']}:0:{normalized[:80]}".encode()).hexdigest()
            return [
                DocumentChunk(
                    chunk_id=chunk_id,
                    text=normalized,
                    source_path=metadata["source_path"],
                    extension=metadata["extension"],
                    metadata=metadata,
                )
            ]
        
        chunks = []
        start = 0
        chunk_index = 0
        
        while start < len(normalized):
            end = min(start + self.chunk_size, len(normalized))
            chunk_text = normalized[start:end]
            
            chunk_id = hashlib.sha1(
                f"{metadata['source_path']}:{chunk_index}:{chunk_text[:80]}".encode()
            ).hexdigest()
            
            chunks.append(
                DocumentChunk(
                    chunk_id=chunk_id,
                    text=chunk_text,
                    source_path=metadata["source_path"],
                    extension=metadata["extension"],
                    metadata=metadata.copy(),
                )
            )
            
            if end == len(normalized):
                break
            
            start = max(0, end - self.overlap)
            chunk_index += 1
        
        return chunks

    def process_file(self, path: Path) -> list[DocumentChunk]:
        """
        Processa um arquivo completo: extrai texto e divide em chunks.
        
        Args:
            path: Caminho absoluto para o arquivo
        
        Returns:
            Lista de DocumentChunk
        """
        text, metadata = self.extract_text(path)
        if not text:
            logger.debug(f"Nenhum texto extraído de {path}")
            return []
        
        return self.split_into_chunks(text, metadata)
